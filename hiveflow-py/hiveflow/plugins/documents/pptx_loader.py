"""PPTX document loader using python-pptx."""

from pathlib import Path

from hiveflow.plugins.documents import Document, DocumentLoaderPlugin


class PptxLoader(DocumentLoaderPlugin):
    """Loader for PPTX files using python-pptx."""

    @property
    def plugin_id(self) -> str:
        return "pptx"

    @property
    def description(self) -> str:
        return "PPTX file loader using python-pptx"

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    async def load(self, file_path: str | Path) -> Document:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required for PPTX support. Install with: pip install python-pptx"
            ) from exc

        path = Path(file_path)
        prs = Presentation(str(path))
        slides: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {i}]"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_parts.append(text)
            # Include speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"Speaker notes: {notes}")
            slides.append("\n".join(slide_parts))

        content = "\n\n".join(slides)
        return Document(
            content=content,
            source=str(path),
            metadata={
                "filename": path.name,
                "size_bytes": path.stat().st_size,
                "slide_count": len(slides),
            },
        )
